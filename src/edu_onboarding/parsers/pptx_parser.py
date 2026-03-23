"""
PPTX Parser - Converts PowerPoint files to Canvas Pages (HTML).

Uses python-pptx to extract text, lists, and structure from slides.
"""

from pathlib import Path
from typing import List, Optional
from pptx import Presentation

from ..models.canvas_models import CanvasPage, WorkflowState
from ..models.migration_report import MigrationError, ErrorSeverity
from ..utils.html_utils import clean_html


class PptxParser:
    """
    Parses PowerPoint (.pptx) files and converts them to CanvasPage objects.
    """
    
    def __init__(self, course_directory: Path):
        self.course_directory = course_directory
        self.errors: List[MigrationError] = []
    
    def parse_pptx(self, file_path: Path, identifier: str = None) -> Optional[CanvasPage]:
        """
        Parse a PPTX file and convert to CanvasPage.
        
        Args:
            file_path: Path to .pptx file
            identifier: Optional identifier (defaults to filename key)
            
        Returns:
            CanvasPage or None
        """
        if not file_path.exists():
            return None
            
        try:
            prs = Presentation(file_path)
            
            # Extract content
            html_content = self._convert_presentation_to_html(prs)
            
            title = file_path.stem.replace('_', ' ').title()
            
            page = CanvasPage(
                title=title,
                identifier=identifier or f"pptx_{file_path.stem}",
                body=html_content,
                workflow_state=WorkflowState.ACTIVE,
                source_file=str(file_path)
            )
            
            return page
            
        except Exception as e:
            self.errors.append(MigrationError(
                severity=ErrorSeverity.WARNING,
                error_type="PPTX_PARSE_ERROR",
                message=f"Failed to parse PPTX: {str(e)}",
                file_path=str(file_path),
                suggested_action="Check if file is valid PowerPoint"
            ))
            return None
    
    def _convert_presentation_to_html(self, prs) -> str:
        """Convert presentation to HTML string"""
        html_parts = []
        
        # Add minimal style
        html_parts.append('<div class="ppt-presentation">')
        
        for i, slide in enumerate(prs.slides):
            html_parts.append(f'<div class="ppt-slide" id="slide-{i+1}" style="margin-bottom: 30px; border: 1px solid #eee; padding: 20px;">')
            
            # Slide Title
            if slide.shapes.title:
                title_text = slide.shapes.title.text
                if title_text.strip():
                    html_parts.append(f'<h2>{clean_html(title_text)}</h2>')
            
            # Slide Content
            content_parts = []
            for shape in slide.shapes:
                if shape == slide.shapes.title:
                    continue
                
                if hasattr(shape, "text") and shape.text.strip():
                    text = clean_html(shape.text)
                    # Convert bullets if possible (simple heuristic)
                    if '\n' in text:
                        # Split by lines and wrap in p or li
                        items = [line.strip() for line in text.split('\n') if line.strip()]
                        content_parts.append('<ul>')
                        for item in items:
                            content_parts.append(f'<li>{item}</li>')
                        content_parts.append('</ul>')
                    else:
                        content_parts.append(f'<p>{text}</p>')
            
            if content_parts:
                html_parts.extend(content_parts)
            
            # Notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text
                if notes.strip():
                    html_parts.append(f'<div class="ppt-notes" style="background: #f9f9f9; padding: 10px; margin-top: 10px; font-size: 0.9em; color: #666;"><strong>Notes:</strong> {clean_html(notes)}</div>')
            
            html_parts.append('</div>')
            html_parts.append('<hr>')
        
        html_parts.append('</div>')
        
        # Add tracking script - Enables LMS progress tracking
        # Emits: PRESENTATION_INIT, SLIDE_VIEWED, PRESENTATION_COMPLETE
        html_parts.append('''
<script>
document.addEventListener('DOMContentLoaded', function() {
    const slides = document.querySelectorAll('.ppt-slide');
    const totalSlides = slides.length;
    const viewedSlides = new Set();
    
    // Notify LMS of initial state
    if (window.parent) {
        window.parent.postMessage({
            type: 'PRESENTATION_INIT',
            slideCount: totalSlides,
            source: 'tutor_lms_pptx'
        }, '*');
    }
    
    const observer = new IntersectionObserver((entries) => {
        entries.forEach(entry => {
            if (entry.isIntersecting) {
                const slideId = entry.target.id;
                const slideIndex = parseInt(slideId.replace('slide-', ''));
                
                if (!viewedSlides.has(slideId)) {
                    viewedSlides.add(slideId);
                    
                    // Notify LMS of progress
                    if (window.parent) {
                        window.parent.postMessage({
                            type: 'SLIDE_VIEWED',
                            slideId: slideId,
                            slideIndex: slideIndex,
                            progress: Math.round((viewedSlides.size / totalSlides) * 100),
                            source: 'tutor_lms_pptx'
                        }, '*');
                        
                        // Check completion
                        if (viewedSlides.size === totalSlides) {
                            window.parent.postMessage({
                                type: 'PRESENTATION_COMPLETE',
                                totalSlides: totalSlides,
                                source: 'tutor_lms_pptx'
                            }, '*');
                        }
                    }
                }
            }
        });
    }, { threshold: 0.5 }); // 50% visibility required
    
    slides.forEach(slide => observer.observe(slide));
});
</script>
''')
        
        # Add tracking script
        html_parts.append('''
<script>
document.addEventListener('DOMContentLoaded', function() {
    const slides = document.querySelectorAll('.ppt-slide');
    const totalSlides = slides.length;
    const viewedSlides = new Set();
    
    // Notify LMS of initial state
    window.parent.postMessage({
        type: 'PRESENTATION_INIT',
        slideCount: totalSlides
    }, '*');
    
    const observer = new IntersectionObserver((entries) => {
        entries.forEach(entry => {
            if (entry.isIntersecting) {
                const slideId = entry.target.id;
                const slideIndex = parseInt(slideId.replace('slide-', ''));
                
                if (!viewedSlides.has(slideId)) {
                    viewedSlides.add(slideId);
                    
                    // Notify LMS of progress
                    window.parent.postMessage({
                        type: 'SLIDE_VIEWED',
                        slideId: slideId,
                        slideIndex: slideIndex,
                        progress: Math.round((viewedSlides.size / totalSlides) * 100)
                    }, '*');
                    
                    // Check completion
                    if (viewedSlides.size === totalSlides) {
                        window.parent.postMessage({
                            type: 'PRESENTATION_COMPLETE',
                            totalSlides: totalSlides
                        }, '*');
                    }
                }
            }
        });
    }, { threshold: 0.5 }); // 50% visibility required
    
    slides.forEach(slide => observer.observe(slide));
});
</script>
''')
        
        return '\n'.join(html_parts)
