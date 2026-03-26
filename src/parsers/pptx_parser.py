"""
PPTX Parser - Converts PowerPoint files to Canvas Pages (HTML).

Uses python-pptx to extract text, lists, and structure from slides.
Also auto-generates a cover thumbnail from the first slide so the
ingestion report no longer flags these decks as missing thumbnails.
"""

from pathlib import Path
from typing import List, Optional
from pptx import Presentation

from models.canvas_models import CanvasPage, WorkflowState
from models.migration_report import MigrationError, ErrorSeverity
from utils.html_utils import clean_html


class PptxParser:
    """
    Parses PowerPoint (.pptx) files and converts them to CanvasPage objects.
    Also extracts the first slide as a cover thumbnail (cover_<stem>.png)
    saved alongside the source file, eliminating missing-thumbnail warnings.
    """

    def __init__(self, course_directory: Path):
        self.course_directory = course_directory
        self.errors: List[MigrationError] = []

    def parse_pptx(self, file_path: Path, identifier: str = None) -> Optional[CanvasPage]:
        """
        Parse a PPTX file and convert to CanvasPage.
        Also auto-generates a cover thumbnail from the first slide.
        """
        if not file_path.exists():
            return None

        try:
            prs = Presentation(file_path)

            # Auto-generate cover thumbnail (best-effort, non-fatal)
            self._extract_cover_thumbnail(prs, file_path)

            html_content = self._convert_presentation_to_html(prs)
            title = file_path.stem.replace('_', ' ').title()

            return CanvasPage(
                title=title,
                identifier=identifier or f"pptx_{file_path.stem}",
                body=html_content,
                workflow_state=WorkflowState.ACTIVE,
                source_file=str(file_path)
            )

        except Exception as e:
            self.errors.append(MigrationError(
                severity=ErrorSeverity.WARNING,
                error_type="PPTX_PARSE_ERROR",
                message=f"Failed to parse PPTX: {str(e)}",
                file_path=str(file_path),
                suggested_action="Check if file is valid PowerPoint"
            ))
            return None

    def _extract_cover_thumbnail(self, prs, file_path: Path) -> Optional[Path]:
        """
        Render the first slide to a PNG thumbnail using Pillow.
        Saves as cover_<stem>.png next to the source file.
        Falls back gracefully if Pillow is not installed.
        """
        cover_path = file_path.parent / f"cover_{file_path.stem}.png"
        if cover_path.exists():
            return cover_path  # already generated, skip

        try:
            from PIL import Image, ImageDraw, ImageFont

            if not prs.slides:
                return None

            slide = prs.slides[0]

            # Convert EMU dimensions to pixels at 96 dpi
            emu_per_inch = 914400
            dpi = 96
            width_px  = max(int(prs.slide_width  / emu_per_inch * dpi), 800)
            height_px = max(int(prs.slide_height / emu_per_inch * dpi), 600)

            img = Image.new("RGB", (width_px, height_px), color=(255, 255, 255))
            draw = ImageDraw.Draw(img)

            # Coloured header bar
            draw.rectangle([0, 0, width_px, height_px // 5], fill=(30, 80, 160))

            # Load fonts (fall back to default if system fonts unavailable)
            try:
                title_font = ImageFont.truetype("arial.ttf", size=36)
                body_font  = ImageFont.truetype("arial.ttf", size=20)
            except Exception:
                title_font = ImageFont.load_default()
                body_font  = title_font

            # Draw slide title in header bar
            title_text = ""
            if slide.shapes.title and slide.shapes.title.text.strip():
                title_text = slide.shapes.title.text.strip()
            if title_text:
                draw.text((40, height_px // 10), title_text, fill=(255, 255, 255), font=title_font)

            # Draw first few body lines below the header
            y = height_px // 5 + 30
            for shape in slide.shapes:
                if shape == slide.shapes.title:
                    continue
                if hasattr(shape, "text") and shape.text.strip():
                    for line in shape.text.strip().splitlines()[:6]:
                        if line.strip():
                            draw.text((40, y), line.strip()[:80], fill=(50, 50, 50), font=body_font)
                            y += 30
                        if y > height_px - 40:
                            break
                if y > height_px - 40:
                    break

            img.save(str(cover_path), "PNG", optimize=True)
            return cover_path

        except ImportError:
            # Pillow not installed — skip silently
            return None
        except Exception as e:
            self.errors.append(MigrationError(
                severity=ErrorSeverity.WARNING,
                error_type="PPTX_THUMBNAIL_ERROR",
                message=f"Could not generate thumbnail for {file_path.name}: {e}",
                file_path=str(file_path),
                suggested_action="Install Pillow: pip install Pillow"
            ))
            return None

    def _convert_presentation_to_html(self, prs) -> str:
        """Convert presentation slides to a single HTML string."""
        html_parts = ['<div class="ppt-presentation">']

        for i, slide in enumerate(prs.slides):
            html_parts.append(
                f'<div class="ppt-slide" id="slide-{i+1}" '
                f'style="margin-bottom:30px;border:1px solid #eee;padding:20px;">'
            )

            # Title
            if slide.shapes.title and slide.shapes.title.text.strip():
                html_parts.append(f'<h2>{clean_html(slide.shapes.title.text)}</h2>')

            # Body shapes
            for shape in slide.shapes:
                if shape == slide.shapes.title:
                    continue
                if not (hasattr(shape, "text") and shape.text.strip()):
                    continue
                text = clean_html(shape.text)
                if '\n' in text:
                    items = [l.strip() for l in text.split('\n') if l.strip()]
                    html_parts.append('<ul>')
                    html_parts.extend(f'<li>{item}</li>' for item in items)
                    html_parts.append('</ul>')
                else:
                    html_parts.append(f'<p>{text}</p>')

            # Speaker notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    html_parts.append(
                        f'<div class="ppt-notes" style="background:#f9f9f9;padding:10px;'
                        f'margin-top:10px;font-size:0.9em;color:#666;">'
                        f'<strong>Notes:</strong> {clean_html(notes)}</div>'
                    )

            html_parts.append('</div><hr>')

        html_parts.append('</div>')

        # Single, safe tracking script — works in both iframe and standalone contexts
        html_parts.append('''<script>
document.addEventListener('DOMContentLoaded', function() {
    var slides = document.querySelectorAll('.ppt-slide');
    var totalSlides = slides.length;
    var viewedSlides = new Set();

    function postToParent(msg) {
        if (window.parent && window.parent !== window) {
            window.parent.postMessage(msg, '*');
        }
    }

    postToParent({ type: 'PRESENTATION_INIT', slideCount: totalSlides, source: 'pptx' });

    var observer = new IntersectionObserver(function(entries) {
        entries.forEach(function(entry) {
            if (!entry.isIntersecting) return;
            var id = entry.target.id;
            if (viewedSlides.has(id)) return;
            viewedSlides.add(id);
            var idx = parseInt(id.replace('slide-', ''));
            var progress = Math.round((viewedSlides.size / totalSlides) * 100);
            postToParent({ type: 'SLIDE_VIEWED', slideId: id, slideIndex: idx, progress: progress, source: 'pptx' });
            if (viewedSlides.size === totalSlides) {
                postToParent({ type: 'PRESENTATION_COMPLETE', totalSlides: totalSlides, source: 'pptx' });
            }
        });
    }, { threshold: 0.5 });

    slides.forEach(function(s) { observer.observe(s); });
});
</script>''')

        return '\n'.join(html_parts)
